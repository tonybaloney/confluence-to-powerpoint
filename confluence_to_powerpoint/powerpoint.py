from pptx import Presentation

def make_presentation(target_name):
    prs = Presentation()
    bullet_slide_layout = prs.slide_layouts[1]
    
    slide = prs.slides.add_slide(bullet_slide_layout)
    shapes = slide.shapes
    
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]
    
    title_shape.text = 'Adding a Bullet Slide'
    
    tf = body_shape.text_frame
    tf.text = 'Find the bullet slide layout'
    
    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.text for first bullet'
    p.level = 1
    
    p = tf.add_paragraph()
    p.text = 'Use _TextFrame.add_paragraph() for subsequent bullets'
    p.level = 2
    
    prs.save(target_name)